import io
import uuid
from typing import Any, List
from langchain.tools import tool
from pptx import Presentation
from pptx.util import Inches

from services.files import FileStore
from services.limits import MAX_PPTX_BULLETS_PER_SLIDE, MAX_PPTX_SLIDES
from services.storage import StorageError
from tools.gating import claim_generation_slot


@tool
def create_pptx(topic: str, content: str) -> str:
    """Create a PowerPoint presentation (.pptx file).

    Use ONLY when the user explicitly asks for a presentation, slides,
    PowerPoint, or visual deck. Do NOT use for general questions,
    text-only responses, or "tell me about" without mentioning slides.

    Args:
        topic: Presentation title shown on the title slide (concise).
        content: Slide text separated by double newlines. Each block's first
            line is the slide title, remaining lines are bullet points.

    Returns:
        Filename of the saved presentation (plus its download ID).
    """
    try:
        # Gate BEFORE expensive work: quota checks must precede generation
        # so exhausted users cannot trigger repeated builds.
        user_id, denied = claim_generation_slot("create_pptx")
        if denied is not None:
            return denied
        # Slide cap BEFORE expensive work: count derived slides first so
        # enormous requests are truncated, never fully built.
        blocks: list[str] = [b for b in content.strip().split("\n\n") if b.strip()]
        total_slides = 1 + len(blocks)
        overflow_note = ""
        if total_slides > MAX_PPTX_SLIDES:
            blocks = blocks[: MAX_PPTX_SLIDES - 1]
            overflow_note = (
                f" [Note: limited to the first {MAX_PPTX_SLIDES} slides "
                f"({total_slides} requested).]"
            )
        prs: Presentation = Presentation()

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic

        slides_content: list[str] = blocks
        dropped_bullets = 0
        for slide_text in slides_content:
            if not slide_text.strip():
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            lines: list[str] = slide_text.strip().split("\n")
            slide.shapes.title.text = lines[0].lstrip("- ").strip() if lines else "Slide"
            bullets: list[str] = [ln.lstrip("- ").strip() for ln in lines[1:] if ln.strip()]
            if len(bullets) > MAX_PPTX_BULLETS_PER_SLIDE:
                dropped_bullets += len(bullets) - MAX_PPTX_BULLETS_PER_SLIDE
                bullets = bullets[:MAX_PPTX_BULLETS_PER_SLIDE]
            if not bullets:
                continue
            try:
                body = slide.placeholders[1].text_frame
                body.clear()
                body.paragraphs[0].text = bullets[0]
                for line in bullets[1:]:
                    p = body.add_paragraph()
                    p.text = line
                    p.level = 0
            except (IndexError, KeyError, AttributeError):
                # Fallback: add a textbox if the layout has no body placeholder.
                left = top = Inches(1)
                width = Inches(8)
                height = Inches(4)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.text = bullets[0]
                for line in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = f"\u2022 {line}"
                    p.level = 0

        filename: str = f"pptx_{uuid.uuid4().hex[:8]}.pptx"
        buf = io.BytesIO()
        prs.save(buf)
        data: bytes = buf.getvalue()

        try:
            meta = FileStore(user_id).register_output(filename, data, "pptx")
            if dropped_bullets:
                overflow_note += (
                    f" [{dropped_bullets} bullet lines dropped "
                    f"(max {MAX_PPTX_BULLETS_PER_SLIDE} per slide).]"
                )
            return (
                f"Presentation saved as {meta.display_name} "
                f"(file ID: {meta.id}){overflow_note}"
            )
        except StorageError as e:
            return f"STATUS=FAILED tool=create_pptx: {e}"
    except Exception as e:
        return f"STATUS=FAILED tool=create_pptx: {str(e)}"


_PPTX_MAX_SLIDES: int = 20
_PPTX_MAX_BULLETS: int = 7
_PPTX_MAX_CHARS_PER_BULLET: int = 160
_PPTX_MAX_TITLE: int = 80
_PPTX_MAX_TABLE_ROWS: int = 12
_PPTX_MAX_TABLE_COLS: int = 6
_PPTX_SLIDE_TYPES = ("title", "section", "bullets", "two_column", "table")


def _pptx_set_textbox(slide: Any, left: Any, top: Any, width: Any, height: Any, title: str, lines: List[str]) -> None:
    """Add a titled textbox; never raises out (best effort)."""
    from pptx.util import Pt

    try:
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        first = tf.paragraphs[0]
        first.text = title[:_PPTX_MAX_TITLE]
        first.runs[0].font.bold = True
        first.runs[0].font.size = Pt(20)
        for line in lines:
            p = tf.add_paragraph()
            p.text = f"• {line[:_PPTX_MAX_CHARS_PER_BULLET]}"
            p.level = 0
            if p.runs:
                p.runs[0].font.size = Pt(14)
    except Exception:
        pass


def _pptx_add_notes(slide: Any, notes: str) -> None:
    """Attach speaker notes; ignored when unsupported."""
    if not notes or not notes.strip():
        return
    try:
        slide.notes_slide.placeholders[1].text = notes.strip()[:2000]
    except Exception:
        pass


def _pptx_add_number(slide: Any, number: int) -> None:
    """Stamp a small slide number bottom-right; ignored when unsupported."""
    try:
        from pptx.util import Pt

        box = slide.shapes.add_textbox(Inches(9.0), Inches(6.8), Inches(1.0), Inches(0.4))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = str(number)
        run.font.size = Pt(10)
    except Exception:
        pass


@tool
def build_presentation(spec_json: str) -> str:
    """Build a structured, validated PowerPoint deck from a JSON spec.

    Prefer this over create_pptx when a designed deck is needed. The spec
    is a JSON object:
      {"title": "Deck title", "subtitle": "optional",
       "slides": [{"type": "bullets", "title": "...", "bullets": [...],
                   "notes": "..."}, ...]}
    Slide types: title, section, bullets, two_column, table.
    - bullets: {"title": str, "bullets": [str, ...]}
    - two_column: {"title": str, "left_title": str, "left": [...],
      "right_title": str, "right": [...]}
    - table: {"title": str, "table": {"headers": [...], "rows": [[...]]}}
    - title/section: {"title": str, "subtitle": str}
    Long bullet lists auto-split into "(cont.)" slides; oversized text is
    truncated with a report instead of overflowing silently. Max 20 slides.

    Args:
        spec_json: The JSON specification string.

    Returns:
        Summary with filename (plus download ID), slide count, and any
        overflow report — or a STATUS= error marker.
    """
    import json as _json

    try:
        spec = _json.loads(spec_json)
    except Exception:
        return "STATUS=INVALID tool=build_presentation: spec_json is not valid JSON."
    if not isinstance(spec, dict) or not isinstance(spec.get("slides"), list):
        return "STATUS=INVALID tool=build_presentation: spec needs a 'slides' list."
    title = str(spec.get("title", "Presentation") or "Presentation")[:_PPTX_MAX_TITLE]
    subtitle = str(spec.get("subtitle", ""))[:120]
    raw_slides = spec["slides"][:_PPTX_MAX_SLIDES]
    if not raw_slides:
        return "STATUS=INVALID tool=build_presentation: no slides in spec."
    for i, s in enumerate(raw_slides):
        if not isinstance(s, dict) or s.get("type", "bullets") not in _PPTX_SLIDE_TYPES:
            return (
                f"STATUS=INVALID tool=build_presentation: slide {i + 1} has "
                f"bad type (allowed: {', '.join(_PPTX_SLIDE_TYPES)})."
            )

    # Gate BEFORE expensive work: quota checks must precede generation
    # so exhausted users cannot trigger repeated builds.
    user_id, denied = claim_generation_slot("build_presentation")
    if denied is not None:
        return denied

    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    overflow_notes: List[str] = []
    try:
        prs: Presentation = Presentation()
        built = 0

        def _new_content_slide(slide_title: str) -> Any:
            nonlocal built
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = (slide_title or "Slide")[:_PPTX_MAX_TITLE]
            built += 1
            _pptx_add_number(slide, built + 1)
            return slide

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        try:
            if subtitle and len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = subtitle
        except Exception:
            pass
        built = 1

        for s in raw_slides:
            if built >= MAX_PPTX_SLIDES:
                overflow_notes.append(
                    f"Stopped at {MAX_PPTX_SLIDES} slides (limit)."
                )
                break
            stype = s.get("type", "bullets")
            stitle = str(s.get("title", "") or "Slide")[:_PPTX_MAX_TITLE]
            notes = str(s.get("notes", ""))

            if stype == "title":
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = stitle
                try:
                    if s.get("subtitle") and len(slide.placeholders) > 1:
                        slide.placeholders[1].text = str(s["subtitle"])[:120]
                except Exception:
                    pass
                built += 1
                _pptx_add_notes(slide, notes)

            elif stype == "section":
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_set_textbox(slide, Inches(0.8), Inches(2.4), Inches(8.4), Inches(2.4),
                                  stitle, [str(s.get("subtitle", ""))[:120]] if s.get("subtitle") else [])
                try:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    run.font.color.rgb = RGBColor(0x63, 0x66, 0xF1)
                                    run.font.bold = True
                except Exception:
                    pass
                built += 1
                _pptx_add_number(slide, built)
                _pptx_add_notes(slide, notes)

            elif stype == "two_column":
                left = [str(x).strip()[:_PPTX_MAX_CHARS_PER_BULLET] for x in (s.get("left") or []) if str(x).strip()]
                right = [str(x).strip()[:_PPTX_MAX_CHARS_PER_BULLET] for x in (s.get("right") or []) if str(x).strip()]
                slide = _new_content_slide(stitle)
                try:
                    slide.placeholders[1].text_frame.clear()
                except Exception:
                    pass
                _pptx_set_textbox(slide, Inches(0.5), Inches(1.6), Inches(4.4), Inches(4.6),
                                  str(s.get("left_title", ""))[:60], left)
                _pptx_set_textbox(slide, Inches(5.1), Inches(1.6), Inches(4.4), Inches(4.6),
                                  str(s.get("right_title", ""))[:60], right)
                _pptx_add_notes(slide, notes)

            elif stype == "table":
                table_spec = s.get("table") or {}
                headers = [str(h)[:40] for h in (table_spec.get("headers") or [])[:_PPTX_MAX_TABLE_COLS]]
                rows = [[str(c)[:60] for c in (r or [])[:_PPTX_MAX_TABLE_COLS]]
                        for r in (table_spec.get("rows") or [])[:_PPTX_MAX_TABLE_ROWS]]
                if not headers or not rows:
                    return (
                        "STATUS=INVALID tool=build_presentation: table slides need "
                        "headers and at least one row."
                    )
                slide = _new_content_slide(stitle)
                try:
                    slide.placeholders[1].text_frame.clear()
                except Exception:
                    pass
                try:
                    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(1.6), Inches(9.0), Inches(4.0))
                    table = shape.table
                    for j, header in enumerate(headers):
                        table.cell(0, j).text = header
                    for i, row in enumerate(rows):
                        for j in range(len(headers)):
                            table.cell(i + 1, j).text = row[j] if j < len(row) else ""
                except Exception as e:
                    return f"STATUS=FAILED tool=build_presentation: table render failed ({str(e)[:120]})."
                _pptx_add_notes(slide, notes)

            else:  # bullets (default)
                bullets = [str(x).strip()[:_PPTX_MAX_CHARS_PER_BULLET]
                           for x in (s.get("bullets") or []) if str(x).strip()]
                if not bullets:
                    overflow_notes.append(f"Slide '{stitle}' had no bullets; skipped.")
                    continue
                chunks = [bullets[i:i + _PPTX_MAX_BULLETS]
                          for i in range(0, len(bullets), _PPTX_MAX_BULLETS)]
                # Total slide cap: never expand past the limit.
                allowed = chunks[: max(0, MAX_PPTX_SLIDES - built)]
                if len(allowed) < len(chunks):
                    overflow_notes.append(
                        f"Slide '{stitle}' truncated at {MAX_PPTX_SLIDES} slides (limit)."
                    )
                for k, chunk in enumerate(allowed):
                    name = stitle if k == 0 else f"{stitle} (cont.)"
                    slide = _new_content_slide(name)
                    try:
                        body = slide.placeholders[1].text_frame
                        body.clear()
                        body.paragraphs[0].text = f"• {chunk[0]}"
                        for line in chunk[1:]:
                            p = body.add_paragraph()
                            p.text = f"• {line}"
                            p.level = 0
                    except (IndexError, KeyError, AttributeError):
                        _pptx_set_textbox(slide, Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.6), name, chunk)
                    _pptx_add_notes(slide, notes if k == 0 else "")
                if len(allowed) > 1:
                    overflow_notes.append(
                        f"Slide '{stitle}' split into {len(allowed)} slides (content overflow)."
                    )

        if built < 2:
            return "STATUS=FAILED tool=build_presentation: no usable slides were built."

        filename: str = f"pptx_{uuid.uuid4().hex[:8]}.pptx"
        buf = io.BytesIO()
        prs.save(buf)
        data: bytes = buf.getvalue()

        # Quality control: reopen and verify structure.
        try:
            check = Presentation(io.BytesIO(data))
            if len(check.slides) != built:
                return (
                    "STATUS=FAILED tool=build_presentation: validation failed "
                    f"(expected {built} slides, found {len(check.slides)})."
                )
            empty_titles = 0
            for slide in check.slides:
                has_title = False
                try:
                    if slide.shapes.title and (slide.shapes.title.text or "").strip():
                        has_title = True
                except Exception:
                    pass
                if not has_title:
                    try:
                        for shape in slide.shapes:
                            if shape.has_text_frame and (shape.text or "").strip():
                                has_title = True
                                break
                    except Exception:
                        pass
                if not has_title:
                    empty_titles += 1
            if empty_titles:
                return (
                    "STATUS=FAILED tool=build_presentation: validation failed "
                    f"({empty_titles} slides have no readable content)."
                )
        except Exception as e:
            return f"STATUS=FAILED tool=build_presentation: output unreadable ({str(e)[:120]})."

        summary = f"Presentation saved as {filename} with {built} slides."
        if overflow_notes:
            summary += " " + " ".join(overflow_notes)
        try:
            meta = FileStore(user_id).register_output(filename, data, "pptx")
            return f"{summary} (file ID: {meta.id})"
        except StorageError as e:
            return f"STATUS=FAILED tool=build_presentation: {e}"
    except Exception as e:
        return f"STATUS=FAILED tool=build_presentation: {str(e)[:200]}"
